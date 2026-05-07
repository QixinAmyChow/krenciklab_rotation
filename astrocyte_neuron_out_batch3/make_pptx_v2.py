"""
Reorganized AQuA2 summary — 4 blocks by cell culture type + faceted ±treatment curves.
Block 1: ATP baseline (3.25.26)
Block 2: Astrocyte only ±ABO (3.31.26, 4.6.26)
Block 3: Astro+Neuron coculture ±ABO (4.1.26, 4.3.26)
Block 4: Organoid/Astroid ±ABO (4.8.26, 4.9.26, 4.10.26)
"""

import os
import numpy as np
import pandas as pd
import matplotlib
matplotlib.use('Agg')
import matplotlib.pyplot as plt
from matplotlib.lines import Line2D
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN

OUT_B1 = '/home/crnlqz/krenciklab/astrocyte+neuron_out'
OUT_B2 = '/home/crnlqz/krenciklab/astrocyte+neuron_out_batch2'
OUT_B3 = '/home/crnlqz/krenciklab/astrocyte+neuron_out_batch3'
FACET_PNG = f'{OUT_B3}/faceted_curves.png'
OUT_PPTX  = f'{OUT_B3}/astro_neuron_AQuA2_v2.pptx'

# ── Block / experiment registry ───────────────────────────────────────────────
BLOCKS = [
    {
        'block_title': 'Block 1 — ATP Baseline',
        'block_color': RGBColor(0x5C, 0x6B, 0xC0),
        'exps': [
            dict(
                id='3.25.26', label='3/25/26 — day28 GCaMP8s Asteroids (spont vs +8ul ATP)',
                png=f'{OUT_B3}/astro_neuron_3_25_26_figure.png',
                xlsx=f'{OUT_B3}/astro_neuron_batch3_AQuA2_summary.xlsx',
                exp_col=None, exp_id=None,
                ctrl='spontaneous', treat='+8ul ATP',
                note='n=2 spont, n=1 ATP (no stats)',
            ),
        ],
    },
    {
        'block_title': 'Block 2 — Astrocyte only  ±ABO',
        'block_color': RGBColor(0x2E, 0x86, 0x78),
        'exps': [
            dict(
                id='3.31.26', label='3/31/26 — Astro only (ctrl vs +amyloid 3d)',
                png=f'{OUT_B1}/astro_neuron_3_31_26_figure.png',
                xlsx=f'{OUT_B1}/astro_neuron_AQuA2_summary.xlsx',
                exp_col='experiment', exp_id='3.31.26',
                ctrl='control', treat='+amyloid', note='',
            ),
            dict(
                id='4.6.26', label='4/6/26 — Astro only (ctrl vs +ABO)',
                png=f'{OUT_B2}/astro_neuron_4_6_26_figure.png',
                xlsx=f'{OUT_B2}/astro_neuron_batch2_AQuA2_summary.xlsx',
                exp_col='experiment', exp_id='4.6.26',
                ctrl='control', treat='+ABO', note='',
            ),
        ],
    },
    {
        'block_title': 'Block 3 — Astro+Neuron Coculture  ±ABO',
        'block_color': RGBColor(0xC0, 0x6B, 0x2B),
        'exps': [
            dict(
                id='4.1.26', label='4/1/26 — Astro+Neuron (ctrl vs +ABO day 4)',
                png=f'{OUT_B1}/astro_neuron_4_1_26_figure.png',
                xlsx=f'{OUT_B1}/astro_neuron_AQuA2_summary.xlsx',
                exp_col='experiment', exp_id='4.1.26',
                ctrl='control', treat='+ABO', note='',
            ),
            dict(
                id='4.3.26', label='4/3/26 — Astro+Neuron (ctrl vs +ABO day 6)',
                png=f'{OUT_B1}/astro_neuron_4_3_26_figure.png',
                xlsx=f'{OUT_B1}/astro_neuron_AQuA2_summary.xlsx',
                exp_col='experiment', exp_id='4.3.26',
                ctrl='control', treat='+ABO', note='',
            ),
        ],
    },
    {
        'block_title': 'Block 4 — Organoid / Astroid  ±ABO',
        'block_color': RGBColor(0x8E, 0x44, 0xAD),
        'exps': [
            dict(
                id='4.8.26', label='4/8/26 — Organoid (ctrl vs +ABO)',
                png=f'{OUT_B2}/astro_neuron_4_8_26_figure.png',
                xlsx=f'{OUT_B2}/astro_neuron_batch2_AQuA2_summary.xlsx',
                exp_col='experiment', exp_id='4.8.26',
                ctrl='control', treat='+ABO',
                note='Very few events (0–3/sample)',
            ),
            dict(
                id='4.9.26', label='4/9/26 — GCaMP8s Asteroids d6 (ctrl vs +amyloid)',
                png=f'{OUT_B2}/astro_neuron_4_9_26_figure.png',
                xlsx=f'{OUT_B2}/astro_neuron_batch2_AQuA2_summary.xlsx',
                exp_col='experiment', exp_id='4.9.26',
                ctrl='control', treat='+amyloid',
                note='High variance: rep3 amyloid=96 events',
            ),
            dict(
                id='4.10.26', label='4/10/26 — GCaMP8s Astroid d7 (ctrl vs +amyloid)',
                png=f'{OUT_B2}/astro_neuron_4_10_26_figure.png',
                xlsx=f'{OUT_B2}/astro_neuron_batch2_AQuA2_summary.xlsx',
                exp_col='experiment', exp_id='4.10.26',
                ctrl='control', treat='+amyloid', note='',
            ),
        ],
    },
]

# colors per date within each block
BLOCK_DATE_COLORS = [
    ['#5C6BC0'],                                      # block1: 1 date
    ['#1B5E20', '#43A047'],                           # block2: 2 dates
    ['#BF360C', '#FF7043'],                           # block3: 2 dates
    ['#6A1B9A', '#AB47BC', '#CE93D8'],                # block4: 3 dates
]


# ── Load all events ──────────────────────────────────────────────────────────
def load_events(exp):
    df = pd.read_excel(exp['xlsx'], sheet_name='All Events')
    if exp['exp_col'] and exp['exp_id']:
        df = df[df[exp['exp_col']] == exp['exp_id']].copy()
    return df


# ── Faceted curves figure ─────────────────────────────────────────────────────
plt.rcParams.update({
    'font.size': 9, 'axes.spines.top': False, 'axes.spines.right': False,
    'axes.linewidth': 0.8, 'xtick.major.width': 0.8, 'ytick.major.width': 0.8,
})

BIN = 30
block_labels = [b['block_title'].split(' — ')[1] for b in BLOCKS]

# Color by condition type: treatment = pink family, control = teal family
C_TREAT = '#E8748A'
C_CTRL  = '#4DADA0'

# Line styles cycle for different days within a block
DAY_STYLES = ['-', '--', ':']

fig, axes = plt.subplots(1, 4, figsize=(16, 4.5), sharey=False)
fig.suptitle('Ca²⁺ event timecourse  ·  color = condition type  ·  line style = experiment day',
             fontsize=10, fontweight='bold', y=1.02)

# ── First pass: compute all curves, track global y max ───────────────────────
all_block_data = []
global_ymax = 0

for block in BLOCKS:
    block_data = []
    for exp in block['exps']:
        df = load_events(exp)
        max_frame = int(df['start_frame'].dropna().max()) + BIN
        bins    = np.arange(0, max_frame + BIN, BIN)
        centers = (bins[:-1] + bins[1:]) / 2
        curves = []
        for cond, is_treat in [(exp['treat'], True), (exp['ctrl'], False)]:
            sub  = df[df['condition'] == cond]
            reps = sorted(sub['rep'].unique(), key=str)
            counts = []
            for rep in reps:
                frames = sub[sub['rep'] == rep]['start_frame'].dropna().values
                c, _ = np.histogram(frames, bins=bins)
                counts.append(c)
            counts = np.array(counts, dtype=float)
            mean_c = counts.mean(axis=0)
            sem_c  = counts.std(axis=0, ddof=1) / np.sqrt(len(counts)) if len(counts) > 1 else np.zeros_like(mean_c)
            global_ymax = max(global_ymax, (mean_c + sem_c).max())
            curves.append((cond, is_treat, centers, mean_c, sem_c))
        block_data.append((exp, curves))
    all_block_data.append(block_data)

y_top = global_ymax * 1.18

# ── Second pass: draw ─────────────────────────────────────────────────────────
for bi, (block, ax, block_data) in enumerate(zip(BLOCKS, axes, all_block_data)):
    legend_handles = []

    for ei, (exp, curves) in enumerate(block_data):
        ls = DAY_STYLES[ei % len(DAY_STYLES)]
        date_short = exp['id']  # e.g. "3.31.26"

        for cond, is_treat, centers, mean_c, sem_c in curves:
            color = C_TREAT if is_treat else C_CTRL
            label = f"{date_short} {cond}"
            ax.errorbar(centers, mean_c, yerr=sem_c, color=color,
                        lw=1.6, ls=ls, marker='o', markersize=2.5,
                        capsize=2.5, elinewidth=0.8, capthick=0.8,
                        label=label)
            legend_handles.append(
                Line2D([0], [0], color=color, lw=1.5, ls=ls, label=label)
            )

    ax.set_ylim(0, y_top)
    ax.set_title(block_labels[bi], fontsize=9, fontweight='bold')
    ax.set_xlabel('Frame', fontsize=8)
    if bi == 0:
        ax.set_ylabel('Ca²⁺ events / bin (mean ± SEM)', fontsize=8)
    else:
        ax.tick_params(labelleft=False)

    ax.legend(handles=legend_handles, frameon=False, fontsize=6.5,
              loc='upper right', handlelength=2.2)

# Shared color legend at figure bottom
fig.legend(
    handles=[Line2D([0],[0], color=C_TREAT, lw=2, label='Treatment (+ABO / +amyloid / +ATP)'),
             Line2D([0],[0], color=C_CTRL,  lw=2, label='Control')],
    loc='lower center', ncol=2, frameon=False, fontsize=8.5,
    bbox_to_anchor=(0.5, -0.06)
)

plt.tight_layout()
fig.savefig(FACET_PNG, dpi=180, bbox_inches='tight')
plt.close(fig)
print(f'Faceted figure saved: {FACET_PNG}')


# ── PowerPoint ────────────────────────────────────────────────────────────────
C_WHITE = RGBColor(0xFF, 0xFF, 0xFF)
C_DARK  = RGBColor(0x2B, 0x2B, 0x2B)
C_GRAY  = RGBColor(0xF0, 0xF0, 0xF0)

prs = Presentation()
prs.slide_width  = Inches(13.33)
prs.slide_height = Inches(7.5)
blank = prs.slide_layouts[6]


def add_title_bar(slide, text, color, note=''):
    bar = slide.shapes.add_shape(1, 0, 0, prs.slide_width, Inches(0.55))
    bar.fill.solid(); bar.fill.fore_color.rgb = color
    bar.line.fill.background()
    tx = slide.shapes.add_textbox(Inches(0.15), Inches(0.06), Inches(10.5), Inches(0.43))
    p = tx.text_frame.paragraphs[0]
    p.text = text
    p.runs[0].font.size = Pt(13); p.runs[0].font.bold = True
    p.runs[0].font.color.rgb = C_WHITE
    if note:
        nt = slide.shapes.add_textbox(Inches(10.5), Inches(0.09), Inches(2.7), Inches(0.38))
        p2 = nt.text_frame.paragraphs[0]
        p2.alignment = PP_ALIGN.RIGHT
        p2.text = '⚠ ' + note
        p2.runs[0].font.size = Pt(7.5)
        p2.runs[0].font.color.rgb = RGBColor(0xFF, 0xCC, 0x44)


def set_cell(cell, text, bold=False, fontsize=8, align=PP_ALIGN.CENTER, bg=None, fg=C_DARK):
    cell.text = ''
    tf = cell.text_frame; tf.word_wrap = False
    p = tf.paragraphs[0]; p.alignment = align
    run = p.add_run(); run.text = str(text)
    run.font.size = Pt(fontsize); run.font.bold = bold
    run.font.color.rgb = fg
    if bg:
        cell.fill.solid(); cell.fill.fore_color.rgb = bg


def add_summary_table(slide, exp, top):
    df = pd.read_excel(exp['xlsx'], sheet_name='Per-Sample Summary')
    if exp['exp_col'] and exp['exp_id']:
        df = df[df[exp['exp_col']] == exp['exp_id']]
    df = df.copy(); df['rep'] = df['rep'].astype(str)
    cols = ['condition', 'rep', 'n_events', 'mean_area', 'mean_dff', 'mean_auc']
    cols = [c for c in cols if c in df.columns]
    hdrs = ['Condition', 'Rep', 'Events', 'Mean Area (px²)', 'Mean Max ΔF/F', 'Mean AUC'][:len(cols)]

    n_rows = len(df) + 1
    tbl = slide.shapes.add_table(n_rows, len(cols),
                                  Inches(0.25), top, Inches(12.83),
                                  Inches(0.28 * n_rows + 0.05)).table
    for ci, w in enumerate([Inches(2.0), Inches(0.7), Inches(0.9), Inches(2.5), Inches(2.3), Inches(2.0)][:len(cols)]):
        tbl.columns[ci].width = w
    for ci, h in enumerate(hdrs):
        set_cell(tbl.cell(0, ci), h, bold=True, fontsize=8,
                 bg=RGBColor(0x3A, 0x3A, 0x3A), fg=C_WHITE)
    for ri, (_, row) in enumerate(df.iterrows(), start=1):
        is_treat = str(row.get('condition', '')) == exp['treat']
        bg = RGBColor(0xFF, 0xEE, 0xF0) if is_treat else RGBColor(0xEE, 0xF7, 0xF6)
        for ci, col in enumerate(cols):
            val = row[col]
            if col in ('mean_area', 'mean_dff', 'mean_auc') and pd.notna(val):
                val = f'{val:.1f}'
            elif pd.isna(val):
                val = '—'
            align = PP_ALIGN.LEFT if ci == 0 else PP_ALIGN.CENTER
            set_cell(tbl.cell(ri, ci), val, fontsize=8, bg=bg, align=align)


# ── Title slide ───────────────────────────────────────────────────────────────
slide = prs.slides.add_slide(blank)
bg = slide.shapes.add_shape(1, 0, 0, prs.slide_width, prs.slide_height)
bg.fill.solid(); bg.fill.fore_color.rgb = RGBColor(0x12, 0x12, 0x22)
bg.line.fill.background()

tx = slide.shapes.add_textbox(Inches(1), Inches(1.8), Inches(11.3), Inches(2))
tf = tx.text_frame; tf.word_wrap = True
p = tf.paragraphs[0]; p.alignment = PP_ALIGN.CENTER
p.text = 'AQuA2 Ca²⁺ Event Analysis'
p.runs[0].font.size = Pt(38); p.runs[0].font.bold = True
p.runs[0].font.color.rgb = C_WHITE

p2 = tf.add_paragraph(); p2.alignment = PP_ALIGN.CENTER
p2.text = 'Astrocyte + Neuron Experiments  ·  3/25/26 – 4/10/26'
p2.runs[0].font.size = Pt(17)
p2.runs[0].font.color.rgb = RGBColor(0xAA, 0xCC, 0xFF)

for yi, line in enumerate([
    'Block 1 · ATP Baseline', 'Block 2 · Astrocyte only ±ABO',
    'Block 3 · Astro+Neuron Coculture ±ABO', 'Block 4 · Organoid / Astroid ±ABO',
]):
    bx = slide.shapes.add_textbox(Inches(2.5 + yi * 2.15), Inches(5.0),
                                   Inches(2.0), Inches(0.5))
    p3 = bx.text_frame.paragraphs[0]; p3.alignment = PP_ALIGN.CENTER
    p3.text = line
    p3.runs[0].font.size = Pt(8.5)
    p3.runs[0].font.color.rgb = [
        RGBColor(0x9F, 0xA8, 0xDA), RGBColor(0x80, 0xCB, 0xC4),
        RGBColor(0xFF, 0xAB, 0x40), RGBColor(0xCE, 0x93, 0xD8),
    ][yi]


# ── Block slides + experiment slides ─────────────────────────────────────────
for block in BLOCKS:
    # Block divider slide
    slide = prs.slides.add_slide(blank)
    bg = slide.shapes.add_shape(1, 0, 0, prs.slide_width, prs.slide_height)
    bg.fill.solid(); bg.fill.fore_color.rgb = block['block_color']
    bg.line.fill.background()
    tx = slide.shapes.add_textbox(Inches(1), Inches(2.8), Inches(11.3), Inches(1.5))
    tf = tx.text_frame; tf.word_wrap = True
    p = tf.paragraphs[0]; p.alignment = PP_ALIGN.CENTER
    p.text = block['block_title']
    p.runs[0].font.size = Pt(36); p.runs[0].font.bold = True
    p.runs[0].font.color.rgb = C_WHITE
    p2 = tf.add_paragraph(); p2.alignment = PP_ALIGN.CENTER
    p2.text = f"{len(block['exps'])} experiment{'s' if len(block['exps'])>1 else ''}"
    p2.runs[0].font.size = Pt(16)
    p2.runs[0].font.color.rgb = RGBColor(0xDD, 0xDD, 0xFF)

    # Individual experiment slides
    for exp in block['exps']:
        slide = prs.slides.add_slide(blank)
        add_title_bar(slide, exp['label'], block['block_color'], exp.get('note', ''))
        # Figure
        slide.shapes.add_picture(exp['png'], Inches(0), Inches(0.58),
                                  width=prs.slide_width, height=Inches(3.6))
        # Table
        add_summary_table(slide, exp, Inches(4.30))
        # Legend
        leg = slide.shapes.add_textbox(Inches(0.25), Inches(7.12), Inches(12), Inches(0.3))
        p_l = leg.text_frame.paragraphs[0]
        p_l.text = f'  ■ Control ({exp["ctrl"]})  ■ Treatment ({exp["treat"]})'
        p_l.runs[0].font.size = Pt(7.5)
        p_l.runs[0].font.color.rgb = RGBColor(0x66, 0x66, 0x66)


# ── Faceted curves slide ──────────────────────────────────────────────────────
slide = prs.slides.add_slide(blank)
bg = slide.shapes.add_shape(1, 0, 0, prs.slide_width, Inches(0.55))
bg.fill.solid(); bg.fill.fore_color.rgb = RGBColor(0x1A, 0x23, 0x4A)
bg.line.fill.background()
tx = slide.shapes.add_textbox(Inches(0.15), Inches(0.06), Inches(13), Inches(0.43))
p = tx.text_frame.paragraphs[0]
p.text = 'Ca²⁺ Event Timecourse — ±Treatment by Cell Culture Type  (each curve = 1 experiment day)'
p.runs[0].font.size = Pt(12); p.runs[0].font.bold = True
p.runs[0].font.color.rgb = C_WHITE

slide.shapes.add_picture(FACET_PNG, Inches(0.1), Inches(0.65),
                          width=Inches(13.13), height=Inches(6.7))

prs.save(OUT_PPTX)
print(f'PPTX saved: {OUT_PPTX}')
print(f'Total slides: {len(prs.slides)}')
