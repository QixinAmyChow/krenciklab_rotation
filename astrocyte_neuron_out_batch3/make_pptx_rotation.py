"""
Rotation talk — Astrocyte Ca²⁺ Signalling in AD-Relevant In Vitro Models
Style: white background, Calibri, plain bold slide titles, thin rule under title.
"""

import os
import pandas as pd
import numpy as np
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.oxml.ns import qn
from lxml import etree

OUT_B1    = '/home/crnlqz/krenciklab/astrocyte+neuron_out'
OUT_B2    = '/home/crnlqz/krenciklab/astrocyte+neuron_out_batch2'
OUT_B3    = '/home/crnlqz/krenciklab/astrocyte+neuron_out_batch3'
FACET_PNG = f'{OUT_B3}/faceted_curves.png'
OUT_PPTX  = f'{OUT_B3}/astro_neuron_rotation_talk.pptx'

FONT     = 'Calibri'
SLIDE_W  = Inches(13.33)
SLIDE_H  = Inches(7.5)

C_WHITE  = RGBColor(0xFF, 0xFF, 0xFF)
C_TITLE  = RGBColor(0x1F, 0x1F, 0x1F)
C_BODY   = RGBColor(0x26, 0x26, 0x26)
C_SUB    = RGBColor(0x59, 0x59, 0x59)
C_RULE   = RGBColor(0xC0, 0xC0, 0xC0)
C_INTER  = RGBColor(0xF0, 0xF4, 0xFF)
C_IBRDR  = RGBColor(0x70, 0x9A, 0xD3)
C_BLUE   = RGBColor(0x1F, 0x3D, 0x7A)
C_TEAL   = RGBColor(0x4D, 0xAD, 0xA0)
C_PINK   = RGBColor(0xE8, 0x74, 0x8A)
C_LGRAY  = RGBColor(0xBB, 0xBB, 0xBB)

prs = Presentation()
prs.slide_width  = SLIDE_W
prs.slide_height = SLIDE_H
blank = prs.slide_layouts[6]


# ── Font helper ───────────────────────────────────────────────────────────────

def _set_font(run, name=FONT):
    rPr = run._r.get_or_add_rPr()
    for child in rPr.findall(qn('a:latin')):
        rPr.remove(child)
    latin = etree.SubElement(rPr, qn('a:latin'))
    latin.set('typeface', name)


# ── Core helpers ──────────────────────────────────────────────────────────────

def white_bg(slide):
    sh = slide.shapes.add_shape(1, 0, 0, SLIDE_W, SLIDE_H)
    sh.fill.solid(); sh.fill.fore_color.rgb = C_WHITE
    sh.line.fill.background()


def rule(slide, top, left=Inches(0.5), width=None, color=C_RULE, thick=Inches(0.018)):
    w = width or (SLIDE_W - Inches(1.0))
    sh = slide.shapes.add_shape(1, left, top, w, thick)
    sh.fill.solid(); sh.fill.fore_color.rgb = color
    sh.line.fill.background()


def tb(slide, text, left, top, width, height,
       size=16, bold=False, italic=False,
       color=None, align=PP_ALIGN.LEFT, wrap=True):
    txb = slide.shapes.add_textbox(left, top, width, height)
    tf  = txb.text_frame; tf.word_wrap = wrap
    p   = tf.paragraphs[0]; p.alignment = align
    run = p.add_run(); run.text = text
    run.font.name = FONT; run.font.size = Pt(size)
    run.font.bold = bold; run.font.italic = italic
    run.font.color.rgb = color or C_BODY
    _set_font(run)
    return txb


# ── Slide title — plain bold, thin rule below ────────────────────────────────

TITLE_TOP   = Inches(0.28)
TITLE_H     = Inches(0.60)
BODY_TOP    = Inches(1.05)
BODY_H      = Inches(6.20)
MARGIN_L    = Inches(0.55)
CONTENT_W   = Inches(12.23)


def slide_title(slide, title, subtitle=None):
    tb(slide, title, MARGIN_L, TITLE_TOP, CONTENT_W, TITLE_H,
       size=28, bold=True, color=C_TITLE)
    if subtitle:
        tb(slide, subtitle,
           MARGIN_L, TITLE_TOP + TITLE_H + Inches(0.07),
           CONTENT_W, Inches(0.33),
           size=13, italic=True, color=C_SUB)


# ── Bullet list ───────────────────────────────────────────────────────────────

def bullet_box(slide, items, left, top, width, height, base_size=15):
    txb = slide.shapes.add_textbox(left, top, width, height)
    tf  = txb.text_frame; tf.word_wrap = True
    first = True
    for text, lvl in items:
        p = tf.paragraphs[0] if first else tf.add_paragraph()
        first = False
        if not text:
            p.space_before = Pt(4); continue
        prefix = ('▸ ' if lvl == 0 else '    – ' if lvl == 1 else '        · ')
        sz  = base_size if lvl == 0 else base_size - 1
        clr = C_BODY if lvl == 0 else C_SUB
        run = p.add_run(); run.text = prefix + text
        run.font.name = FONT; run.font.size = Pt(sz)
        run.font.bold = (lvl == 0); run.font.color.rgb = clr
        _set_font(run)
        p.space_before = Pt(7 if lvl == 0 else 3)


# ── Takeaway box ──────────────────────────────────────────────────────────────

def takeaway(slide, lines, left, top, width, height, base_size=13):
    sh = slide.shapes.add_shape(1, left, top, width, height)
    sh.fill.solid(); sh.fill.fore_color.rgb = C_INTER
    sh.line.color.rgb = C_IBRDR; sh.line.width = 12700
    txb = slide.shapes.add_textbox(
        left + Inches(0.15), top + Inches(0.10),
        width - Inches(0.30), height - Inches(0.20))
    tf  = txb.text_frame; tf.word_wrap = True
    first = True
    for i, line in enumerate(lines):
        p = tf.paragraphs[0] if first else tf.add_paragraph()
        first = False
        run = p.add_run(); run.text = line
        run.font.name = FONT
        run.font.size = Pt(base_size + 1 if i == 0 else base_size)
        run.font.bold = (i == 0)
        run.font.color.rgb = C_BLUE
        _set_font(run)
        p.space_before = Pt(0 if i == 0 else 3)


# ── Figure helper ─────────────────────────────────────────────────────────────

def add_fig(slide, path, left, top, width=None, height=None):
    if not os.path.exists(path):
        ph = slide.shapes.add_shape(1, left, top,
                                     width or Inches(8), height or Inches(4))
        ph.fill.solid(); ph.fill.fore_color.rgb = RGBColor(0xF0, 0xF0, 0xF0)
        ph.line.color.rgb = C_LGRAY
        r = slide.shapes.add_textbox(
            left + Inches(0.1), top + Inches(0.1),
            (width or Inches(8)) - Inches(0.2), Inches(0.3)
        ).text_frame.paragraphs[0].add_run()
        r.text = f'[missing: {os.path.basename(path)}]'
        r.font.color.rgb = C_LGRAY; r.font.size = Pt(10)
        return
    kw = {}
    if width:  kw['width']  = width
    if height: kw['height'] = height
    slide.shapes.add_picture(path, left, top, **kw)


# ── Section divider ───────────────────────────────────────────────────────────

def section_div(slide, title, subtitle=''):
    white_bg(slide)
    rule(slide, Inches(2.8), left=Inches(1.5),
         width=Inches(10.33), color=RGBColor(0x80, 0x80, 0x80))
    txb = slide.shapes.add_textbox(Inches(1.2), Inches(2.2), Inches(10.93), Inches(2.2))
    tf  = txb.text_frame; tf.word_wrap = True
    p = tf.paragraphs[0]; p.alignment = PP_ALIGN.CENTER
    run = p.add_run(); run.text = title
    run.font.name = FONT; run.font.size = Pt(40)
    run.font.bold = True; run.font.color.rgb = C_TITLE
    _set_font(run)
    if subtitle:
        p2 = tf.add_paragraph(); p2.alignment = PP_ALIGN.CENTER
        r2 = p2.add_run(); r2.text = subtitle
        r2.font.name = FONT; r2.font.size = Pt(16)
        r2.font.color.rgb = C_SUB; _set_font(r2)
        p2.space_before = Pt(10)


# ── Data summary table ────────────────────────────────────────────────────────

EXPS = {
    '3.25.26': dict(xlsx=f'{OUT_B3}/astro_neuron_batch3_AQuA2_summary.xlsx',
                    exp_col=None,           exp_id=None,
                    ctrl='spontaneous',     treat='+8ul ATP',
                    png=f'{OUT_B3}/astro_neuron_3_25_26_figure.png'),
    '3.31.26': dict(xlsx=f'{OUT_B1}/astro_neuron_AQuA2_summary.xlsx',
                    exp_col='experiment',   exp_id='3.31.26',
                    ctrl='control',         treat='+amyloid',
                    png=f'{OUT_B1}/astro_neuron_3_31_26_figure.png'),
    '4.1.26':  dict(xlsx=f'{OUT_B1}/astro_neuron_AQuA2_summary.xlsx',
                    exp_col='experiment',   exp_id='4.1.26',
                    ctrl='control',         treat='+ABO',
                    png=f'{OUT_B1}/astro_neuron_4_1_26_figure.png'),
    '4.3.26':  dict(xlsx=f'{OUT_B1}/astro_neuron_AQuA2_summary.xlsx',
                    exp_col='experiment',   exp_id='4.3.26',
                    ctrl='control',         treat='+ABO',
                    png=f'{OUT_B1}/astro_neuron_4_3_26_figure.png'),
    '4.6.26':  dict(xlsx=f'{OUT_B2}/astro_neuron_batch2_AQuA2_summary.xlsx',
                    exp_col='experiment',   exp_id='4.6.26',
                    ctrl='control',         treat='+ABO',
                    png=f'{OUT_B2}/astro_neuron_4_6_26_figure.png'),
    '4.8.26':  dict(xlsx=f'{OUT_B2}/astro_neuron_batch2_AQuA2_summary.xlsx',
                    exp_col='experiment',   exp_id='4.8.26',
                    ctrl='control',         treat='+ABO',
                    png=f'{OUT_B2}/astro_neuron_4_8_26_figure.png'),
    '4.9.26':  dict(xlsx=f'{OUT_B2}/astro_neuron_batch2_AQuA2_summary.xlsx',
                    exp_col='experiment',   exp_id='4.9.26',
                    ctrl='control',         treat='+amyloid',
                    png=f'{OUT_B2}/astro_neuron_4_9_26_figure.png'),
    '4.10.26': dict(xlsx=f'{OUT_B2}/astro_neuron_batch2_AQuA2_summary.xlsx',
                    exp_col='experiment',   exp_id='4.10.26',
                    ctrl='control',         treat='+amyloid',
                    png=f'{OUT_B2}/astro_neuron_4_10_26_figure.png'),
}


def _load_summary(eid):
    e = EXPS[eid]
    try:
        df = pd.read_excel(e['xlsx'], sheet_name='Per-Sample Summary')
        if e['exp_col']:
            df = df[df[e['exp_col']] == e['exp_id']]
        return df.copy()
    except Exception:
        return pd.DataFrame()


def _fmt(v):
    if pd.isna(v): return '—'
    try:   return f'{float(v):.1f}'
    except: return str(v)


def data_table(slide, eid, left, top, width):
    e  = EXPS[eid]
    df = _load_summary(eid)
    if df.empty: return
    cols = [c for c in ['condition', 'rep', 'n_events', 'mean_area', 'mean_dff'] if c in df.columns]
    hdrs = {'condition': 'Condition', 'rep': 'Rep',
            'n_events': 'Events', 'mean_area': 'Area (px²)', 'mean_dff': 'Max ΔF/F'}
    row_h = Inches(0.30)
    tbl   = slide.shapes.add_table(
        len(df) + 1, len(cols), left, top, width, row_h * (len(df) + 1)).table
    cw = {'condition': Inches(2.0), 'rep': Inches(0.5),
          'n_events': Inches(0.8), 'mean_area': Inches(1.1), 'mean_dff': Inches(1.0)}
    for ci, col in enumerate(cols):
        tbl.columns[ci].width = cw.get(col, Inches(1.0))
    for ci, col in enumerate(cols):
        cell = tbl.cell(0, ci)
        cell.text = hdrs[col]
        cell.fill.solid(); cell.fill.fore_color.rgb = RGBColor(0x26, 0x26, 0x26)
        p = cell.text_frame.paragraphs[0]; p.alignment = PP_ALIGN.CENTER
        r = p.runs[0]; r.font.size = Pt(10); r.font.bold = True
        r.font.name = FONT; r.font.color.rgb = C_WHITE
    for ri, (_, row) in enumerate(df.iterrows(), start=1):
        is_t = str(row.get('condition', '')) == e['treat']
        bg   = RGBColor(0xFF, 0xEE, 0xF2) if is_t else RGBColor(0xF0, 0xF8, 0xF6)
        for ci, col in enumerate(cols):
            cell = tbl.cell(ri, ci)
            val  = _fmt(row[col]) if col not in ('condition', 'rep') else str(row[col])
            cell.text = val
            cell.fill.solid(); cell.fill.fore_color.rgb = bg
            p = cell.text_frame.paragraphs[0]
            p.alignment = PP_ALIGN.LEFT if ci == 0 else PP_ALIGN.CENTER
            r = p.runs[0]; r.font.size = Pt(10); r.font.name = FONT


# ── Data slide constants ──────────────────────────────────────────────────────

TONLY_BOT = BODY_TOP      # figure starts right at body top
FIG_W     = Inches(8.6)
FIG_H     = Inches(4.1)
TBL_L     = MARGIN_L + FIG_W + Inches(0.25)
TBL_W     = CONTENT_W - FIG_W - Inches(0.25)
BOX_TOP   = TONLY_BOT + FIG_H + Inches(0.1)
BOX_H     = SLIDE_H - BOX_TOP - Inches(0.12)


def data_slide(title, subtitle, eid, takeaway_lines):
    slide = prs.slides.add_slide(blank)
    white_bg(slide)
    slide_title(slide, title, subtitle=subtitle)
    add_fig(slide, EXPS[eid]['png'], MARGIN_L, TONLY_BOT, width=FIG_W, height=FIG_H)
    data_table(slide, eid, TBL_L, TONLY_BOT, TBL_W)
    takeaway(slide, takeaway_lines, MARGIN_L, BOX_TOP, CONTENT_W, BOX_H, base_size=13)
    return slide


# =============================================================================
# SLIDES
# =============================================================================

# ── 1. Title slide ─────────────────────────────────────────────────────────────
slide = prs.slides.add_slide(blank)
white_bg(slide)

txb = slide.shapes.add_textbox(Inches(0.55), Inches(0.45), Inches(12.23), Inches(1.2))
tf  = txb.text_frame; tf.word_wrap = True
p   = tf.paragraphs[0]
run = p.add_run()
run.text = 'Astrocyte Ca²⁺ Signalling in AD-Relevant In Vitro Models'
run.font.name = FONT; run.font.size = Pt(38); run.font.bold = True
run.font.color.rgb = C_TITLE; _set_font(run)

txb2 = slide.shapes.add_textbox(Inches(0.55), Inches(1.9), Inches(12.23), Inches(0.5))
r2   = txb2.text_frame.paragraphs[0].add_run()
r2.text = 'AQuA2 Calcium Event Analysis  ·  Asteroids  ·  Co-cultures  ·  Organoids'
r2.font.name = FONT; r2.font.size = Pt(18); r2.font.color.rgb = C_SUB; _set_font(r2)

for yi, (label, val) in enumerate([('Student', '[Your Name]'), ('Mentor', '[Mentor Name]'),
                                    ('Lab', 'Krencik Lab'), ('Date', 'May 2026')]):
    row_top = Inches(2.7) + Inches(yi * 0.52)
    tb(slide, label + ':', Inches(0.55), row_top, Inches(1.4), Inches(0.45),
       size=15, bold=True, color=C_BODY)
    tb(slide, val, Inches(2.0), row_top, Inches(5.0), Inches(0.45),
       size=15, color=C_SUB)

tb(slide,
   'Background:  Sanchez-Mico, Calvo-Rodriguez & Bacskai  —  '
   '"Role of dysregulated calcium homeostasis in astrocytes in neurodegenerative disorders"  '
   'Nat Rev Neurosci 2026',
   Inches(0.55), Inches(5.3), Inches(12.23), Inches(0.6),
   size=12, italic=True, color=C_SUB)

tb(slide, 'Rotation Talk', Inches(9.5), Inches(0.45), Inches(3.3), Inches(0.45),
   size=14, color=C_SUB, align=PP_ALIGN.RIGHT)


# ── 2. Background — Astrocyte Ca²⁺ physiology ────────────────────────────────
slide = prs.slides.add_slide(blank)
white_bg(slide)
slide_title(slide, 'Astrocytes: Dynamic Regulators of Brain Homeostasis')
bullet_box(slide, [
    ('Astrocytes are the most abundant glial cell in the CNS', 0),
    ('Neurotransmitter clearance, K⁺ buffering, metabolic coupling, neurovascular regulation', 1),
    ('Form tripartite synapses — integrate and modulate synaptic transmission', 1),
    ('', 0),
    ('Ca²⁺ signalling is the primary "ionic excitability" of astrocytes', 0),
    ('Unlike neurons, astrocytes lack action potentials', 1),
    ('Encode information through spatiotemporally patterned Ca²⁺ signals', 1),
    ('Trigger gliotransmitter release: glutamate, ATP, D-serine → synaptic modulation', 1),
    ('', 0),
    ('Compartmentalized Ca²⁺ architecture  (Sanchez-Mico et al., Nat Rev Neurosci 2026)', 0),
    ('Soma / primary branches  →  slow, IP₃R-ER-dependent transients', 1),
    ('Fine branchlets / perisynaptic leaflets  →  fast, localized microdomains', 1),
    ('Vascular endfeet  →  neurovascular coupling signals', 1),
    ('', 0),
    ('Disruption of astrocyte Ca²⁺ homeostasis  →  synaptic failure, neuroinflammation, circuit dysfunction', 0),
], MARGIN_L, BODY_TOP, CONTENT_W, BODY_H, base_size=15)


# ── 3. Background — AD and Aβ ────────────────────────────────────────────────
slide = prs.slides.add_slide(blank)
white_bg(slide)
slide_title(slide, 'Astrocyte Ca²⁺ Dysregulation in Alzheimer\'s Disease')
bullet_box(slide, [
    ('AD: progressive cognitive decline; Aβ accumulates years before neuronal loss', 0),
    ('Astrocytes respond early — Ca²⁺ changes precede overt plaque deposition', 1),
    ('', 0),
    ('Aβ disrupts astrocyte Ca²⁺ homeostasis through multiple mechanisms:', 0),
    ('↑ Spontaneous Ca²⁺ hyperactivity before plaques appear', 1),
    ('↑ mGluR5 and IP₃R expression  →  sensitized ER Ca²⁺ release', 1),
    ('Loss of STIM1  →  impaired SOCE  →  disrupted synaptic plasticity', 1),
    ('Mitochondrial Ca²⁺ overload  →  ROS, oxidative stress, metabolic failure', 1),
    ('↓ EAAT2  →  glutamate excitotoxicity  ·  Calcineurin / NF-κB  →  neuroinflammation', 1),
    ('', 0),
    ('Network consequence', 0),
    ('Astrocyte Ca²⁺ dysfunction drives early network hyperactivity and neurovascular uncoupling', 1),
    ('These alterations precede neuronal loss — astrocytes are early therapeutic targets', 1),
    ('', 0),
    ('In vitro tool: soluble Aβ oligomers (ABO) — sufficient to recapitulate Ca²⁺ dysregulation', 0),
], MARGIN_L, BODY_TOP, CONTENT_W, BODY_H, base_size=15)


# ── 4. Section divider — Rotation Question ───────────────────────────────────
slide = prs.slides.add_slide(blank)
section_div(slide, 'Rotation Question & Hypothesis')


# ── 5. Rotation Question ──────────────────────────────────────────────────────
slide = prs.slides.add_slide(blank)
white_bg(slide)
slide_title(slide, 'Rotation Question & Hypothesis')

q_top = BODY_TOP
q_sh  = slide.shapes.add_shape(1, MARGIN_L, q_top, CONTENT_W, Inches(1.15))
q_sh.fill.solid(); q_sh.fill.fore_color.rgb = C_INTER
q_sh.line.color.rgb = C_IBRDR; q_sh.line.width = 15000
q_tb = slide.shapes.add_textbox(
    MARGIN_L + Inches(0.2), q_top + Inches(0.12),
    CONTENT_W - Inches(0.4), Inches(0.95))
q_tf = q_tb.text_frame; q_tf.word_wrap = True
q_p  = q_tf.paragraphs[0]; q_p.alignment = PP_ALIGN.CENTER
q_r  = q_p.add_run()
q_r.text = ('Does exposure to soluble amyloid-β oligomers (ABO) alter Ca²⁺ event dynamics '
            'in astrocytes — and does the presence of neurons modulate this response?')
q_r.font.name = FONT; q_r.font.size = Pt(17); q_r.font.bold = True
q_r.font.color.rgb = C_BLUE; _set_font(q_r)

bullet_box(slide, [
    ('Hypothesis', 0),
    ('ABO increases Ca²⁺ event frequency, area, and amplitude in astrocytes', 1),
    ('Astrocyte + neuron co-cultures show amplified or altered responses vs astrocyte-only', 1),
    ('3D organoid models may show attenuated responses due to architectural buffering', 1),
    ('', 0),
    ('Three complementary model systems:', 0),
    ('Asteroid organoids (day 28, GCaMP8s)  —  ATP-validated baseline', 1),
    ('2D monoculture astrocytes ± neurons  —  controlled ±ABO system', 1),
    ('Cerebral organoids / Astroids (day 6–7)  —  3D architecture', 1),
    ('', 0),
    ('Readout: AQuA2 automated Ca²⁺ event detection', 0),
    ('Event count, area (px²), max ΔF/F, AUC  ·  n = 3 replicates per condition', 1),
], MARGIN_L, q_top + Inches(1.25), CONTENT_W, BODY_H - Inches(1.25), base_size=15)


# ── 6. Methods ────────────────────────────────────────────────────────────────
slide = prs.slides.add_slide(blank)
white_bg(slide)
slide_title(slide, 'Experimental Approach')

col_titles = ['GCaMP8s Imaging', 'ABO Treatment', 'AQuA2 Analysis']
col_items  = [
    ['Genetically encoded fast Ca²⁺ indicator',
     'Lentiviral expression in astrocytes',
     'Widefield fluorescence imaging',
     'LIF → TIFF → AQuA2 pipeline'],
    ['Soluble amyloid-β oligomers (ABO)',
     'AD-relevant concentration',
     'Day 4 and Day 6 post-treatment timepoints',
     'Paired control wells in all experiments'],
    ['Automated Ca²⁺ event detection',
     'Spatiotemporal event segmentation',
     'Per-event: area, max ΔF/F, AUC, duration',
     'n = 3 reps/condition  ·  mean ± SEM'],
]
col_w = Inches(3.9)
for xi in range(3):
    lft = MARGIN_L + xi * Inches(4.1)
    tb(slide, col_titles[xi], lft, BODY_TOP, col_w, Inches(0.45),
       size=17, bold=True, color=C_TITLE)
    rule(slide, BODY_TOP + Inches(0.48), left=lft, width=col_w,
         color=RGBColor(0x80, 0x80, 0x80))
    for yi, line in enumerate(col_items[xi]):
        tb(slide, '• ' + line,
           lft, BODY_TOP + Inches(0.6) + Inches(yi * 0.52), col_w, Inches(0.48),
           size=14, color=C_BODY)

rule(slide, BODY_TOP + Inches(2.85), color=C_RULE)
tb(slide, 'Batch overview', MARGIN_L, BODY_TOP + Inches(3.0),
   CONTENT_W, Inches(0.42), size=16, bold=True, color=C_TITLE)
for yi, (b, desc) in enumerate([
    ('Batch 3  (3/25/26)',          'Asteroid day28 — spontaneous + ATP stimulation  [system validation]'),
    ('Batch 1  (3/31/26 – 4/3/26)', 'Astrocyte-only and Astro+Neuron ±ABO  (4 experiments)'),
    ('Batch 2  (4/6/26 – 4/10/26)', 'Astrocyte-only, Organoid, Astroid ±ABO  (4 experiments)'),
]):
    row_top = BODY_TOP + Inches(3.5) + Inches(yi * 0.57)
    tb(slide, b,    MARGIN_L,                row_top, Inches(2.9), Inches(0.5), size=14, bold=True,  color=C_BODY)
    tb(slide, desc, MARGIN_L + Inches(3.05), row_top, Inches(9.3), Inches(0.5), size=14, color=C_SUB)


# ── 7. Section divider — Results ─────────────────────────────────────────────
slide = prs.slides.add_slide(blank)
section_div(slide, 'Results',
            '1 · System Validation    2 · Astrocyte-only ±ABO    '
            '3 · Astro+Neuron ±ABO    4 · Organoid ±ABO')


# ── 8–12. Data slides ─────────────────────────────────────────────────────────
data_slide(
    'Result 1 — System Validation: Spontaneous vs ATP-Evoked Activity',
    '3/25/26  ·  Day 28 GCaMP8s Asteroids  ·  n = 2 spontaneous, n = 1 +ATP',
    '3.25.26',
    ['Key takeaways',
     '• GCaMP8s asteroids at day 28 show robust spontaneous Ca²⁺ events',
     '• ATP stimulation produces a clear increase in event frequency — system is stimulus-responsive',
     '• n = 1 for +ATP precludes statistics; used as proof-of-concept only'])

data_slide(
    'Result 2 — Astrocyte-only: Response to Amyloid-β Oligomers',
    '3/31/26  ·  Astrocyte-only  ·  control vs +amyloid (3-day)  ·  n = 3',
    '3.31.26',
    ['Key takeaways',
     '• First test of amyloid exposure in astrocyte-only cultures',
     '• Isolates the astrocyte-intrinsic component of the ABO response',
     '• Compare event count, area, and max ΔF/F between control and +amyloid groups'])

data_slide(
    'Result 2 (cont.) — Astrocyte-only: ABO Response',
    '4/6/26  ·  Astrocyte-only  ·  control vs +ABO  ·  n = 3',
    '4.6.26',
    ['Key takeaways',
     '• Replication of the astrocyte-only ABO experiment with refined protocol',
     '• Assess consistency of Ca²⁺ response across batches (cf. 3/31/26)',
     '• Do astrocytes alone respond robustly to ABO at this concentration?'])

data_slide(
    'Result 3 — Astro + Neuron Co-culture: ABO Response (Day 4)',
    '4/1/26  ·  Astro+Neuron  ·  control vs +ABO day 4  ·  n = 3',
    '4.1.26',
    ['Key takeaways',
     '• Neurons present — does their activity amplify or buffer ABO-driven Ca²⁺ changes?',
     '• Day 4 post-ABO: captures early-phase Ca²⁺ dysregulation',
     '• Key comparison: co-culture vs astrocyte-only at matched ABO dose'])

data_slide(
    'Result 3 (cont.) — Astro + Neuron Co-culture: ABO Response (Day 6)',
    '4/3/26  ·  Astro+Neuron  ·  control vs +ABO day 6  ·  n = 3',
    '4.3.26',
    ['Key takeaways',
     '• Day 6 post-ABO: does Ca²⁺ dysregulation persist or escalate?',
     '• Two timepoints (day 4 + day 6) track temporal trajectory of dysregulation',
     '• Compare day 4 vs day 6: sustained vs transient response?'])


# ── 13. Result 4 — Organoid / Astroid ────────────────────────────────────────
slide = prs.slides.add_slide(blank)
white_bg(slide)
slide_title(slide, 'Result 4 — Organoid / Astroid: ABO Response in 3D Models',
            subtitle='4/8/26 – 4/10/26  ·  Organoid + GCaMP8s Astroids ±ABO  ·  n = 3 per experiment')

fig_w3 = Inches(4.05)
fig_h3 = Inches(3.6)
for xi, eid in enumerate(['4.8.26', '4.9.26', '4.10.26']):
    lft = MARGIN_L + xi * Inches(4.28)
    add_fig(slide, EXPS[eid]['png'], lft, BODY_TOP, width=fig_w3, height=fig_h3)
    tb(slide, ['4/8/26  Organoid ±ABO',
               '4/9/26  GCaMP8s Astroid d6 ±amyloid  ⚠ high variance',
               '4/10/26  GCaMP8s Astroid d7 ±amyloid'][xi],
       lft, BODY_TOP + fig_h3 + Inches(0.05), fig_w3, Inches(0.4),
       size=12, bold=True, color=C_BODY, align=PP_ALIGN.CENTER)

box_top3 = BODY_TOP + fig_h3 + Inches(0.5)
takeaway(slide, [
    'Key takeaways',
    '• Organoids (4/8/26): very few events (0–3/sample) — 3D architecture may limit excitability or detection',
    '• Astroid d6 (4/9/26): high inter-replicate variance (rep3 = 96 events) — organoid heterogeneity',
    '• Astroid d7 (4/10/26): more consistent; compare to d6 for maturation-state effects',
], MARGIN_L, box_top3, CONTENT_W, SLIDE_H - box_top3 - Inches(0.12), base_size=13)


# ── 14. Faceted timecourse ────────────────────────────────────────────────────
slide = prs.slides.add_slide(blank)
white_bg(slide)
slide_title(slide, 'Ca²⁺ Event Timecourse — All Experiments Overview',
            subtitle='Teal = control  ·  Pink = treatment  ·  Line style = experiment day')
add_fig(slide, FACET_PNG,
        Inches(0.1), BODY_TOP - Inches(0.1),
        width=Inches(13.13), height=SLIDE_H - BODY_TOP)


# ── 15. Section divider — Summary ────────────────────────────────────────────
slide = prs.slides.add_slide(blank)
section_div(slide, 'Summary & Future Directions')


# ── 16. Summary ──────────────────────────────────────────────────────────────
slide = prs.slides.add_slide(blank)
white_bg(slide)
slide_title(slide, 'Summary')
bullet_box(slide, [
    ('System validation  (3/25/26)', 0),
    ('GCaMP8s asteroids at day 28 show robust activity responsive to ATP stimulation', 1),
    ('System is suitable for quantifying treatment-driven Ca²⁺ changes', 1),
    ('', 0),
    ('Astrocyte-only responses to ABO  (3/31/26 · 4/6/26)', 0),
    ('Establishes the astrocyte-intrinsic component of the ABO Ca²⁺ response', 1),
    ('Tests whether ABO recapitulates Aβ-driven Ca²⁺ dysregulation in vitro', 1),
    ('', 0),
    ('Astro + Neuron co-culture  (4/1/26 day 4 · 4/3/26 day 6)', 0),
    ('Addresses whether neurons modulate the astrocyte Ca²⁺ response to ABO', 1),
    ('Two timepoints reveal early vs sustained Ca²⁺ dysregulation', 1),
    ('', 0),
    ('Organoid / Astroid  (4/8/26 – 4/10/26)', 0),
    ('3D models show lower event rates and higher variance — organoid heterogeneity', 1),
    ('More physiologically complex context for evaluating ABO effects', 1),
], MARGIN_L, BODY_TOP, CONTENT_W, BODY_H, base_size=15)


# ── 17. Future Directions ─────────────────────────────────────────────────────
slide = prs.slides.add_slide(blank)
white_bg(slide)
slide_title(slide, 'Future Directions')
bullet_box(slide, [
    ('Increase n for ATP experiments (n = 1)  →  enable statistical comparisons', 0),
    ('', 0),
    ('Dose-response curve for ABO  →  identify effective threshold per model', 0),
    ('', 0),
    ('Subcellular compartment-resolved imaging', 0),
    ('Resolve soma vs branchlet vs endfoot responses — fine-process microdomains most affected early in AD', 1),
    ('', 0),
    ('Two-photon live imaging in 3D organoids  →  improve resolution, reduce variance', 0),
    ('', 0),
    ('Pharmacological rescue experiments', 0),
    ('STIM1 restoration  ·  calcineurin inhibition  ·  RyR blockade — test mechanistic hypotheses', 1),
    ('', 0),
    ('Patient-derived iPSC astrocytes  →  compare sporadic vs familial AD Ca²⁺ phenotypes', 0),
], MARGIN_L, BODY_TOP, CONTENT_W, BODY_H, base_size=15)


# ── 18. Acknowledgements ──────────────────────────────────────────────────────
slide = prs.slides.add_slide(blank)
white_bg(slide)
slide_title(slide, 'Acknowledgements')
bullet_box(slide, [
    ('Krencik Lab', 0),
    ('[Mentor Name]  —  guidance and support throughout the rotation', 1),
    ('Lab members  —  training on cell culture, imaging, and AQuA2 analysis', 1),
    ('', 0),
    ('Key Reference', 0),
    ('Sanchez-Mico MV, Calvo-Rodriguez M & Bacskai BJ.', 1),
    ('"Role of dysregulated calcium homeostasis in astrocytes in neurodegenerative disorders."', 1),
    ('Nature Reviews Neuroscience (2026).  doi: 10.1038/s41583-026-01032-6', 1),
    ('', 0),
    ('Funding / Resources', 0),
    ('[Funding sources, core facility support, etc.]', 1),
], MARGIN_L, BODY_TOP, CONTENT_W, BODY_H, base_size=15)


# =============================================================================
prs.save(OUT_PPTX)
print(f'Saved: {OUT_PPTX}')
print(f'Total slides: {len(prs.slides)}')
