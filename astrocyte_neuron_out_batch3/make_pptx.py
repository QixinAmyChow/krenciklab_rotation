"""
Build summary PowerPoint for all astrocyte+neuron AQuA2 experiments.
One slide per experiment: figure panels + per-sample data table.
"""

import pandas as pd
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt
import os

OUT_B1 = '/home/crnlqz/krenciklab/astrocyte+neuron_out'
OUT_B2 = '/home/crnlqz/krenciklab/astrocyte+neuron_out_batch2'
OUT_B3 = '/home/crnlqz/krenciklab/astrocyte+neuron_out_batch3'

# ── Experiment registry (chronological order) ────────────────────────────────
EXPERIMENTS = [
    dict(
        label='3/25/26 — day28 GCaMP8s Asteroids (spontaneous vs +8ul ATP)',
        png=f'{OUT_B3}/astro_neuron_3_25_26_figure.png',
        xlsx=f'{OUT_B3}/astro_neuron_batch3_AQuA2_summary.xlsx',
        exp_col=None,  # no experiment column in batch3
        exp_id=None,
        treat='+8ul ATP', ctrl='spontaneous',
        note='n=2 spontaneous reps, n=1 ATP rep (no stats)',
    ),
    dict(
        label='3/31/26 — Astro only (control vs +amyloid 3d)',
        png=f'{OUT_B1}/astro_neuron_3_31_26_figure.png',
        xlsx=f'{OUT_B1}/astro_neuron_AQuA2_summary.xlsx',
        exp_col='experiment', exp_id='3.31.26',
        treat='+amyloid', ctrl='control',
        note='',
    ),
    dict(
        label='4/1/26 — Astro+Neuron (control vs +ABO day 4)',
        png=f'{OUT_B1}/astro_neuron_4_1_26_figure.png',
        xlsx=f'{OUT_B1}/astro_neuron_AQuA2_summary.xlsx',
        exp_col='experiment', exp_id='4.1.26',
        treat='+ABO', ctrl='control',
        note='',
    ),
    dict(
        label='4/3/26 — Astro+Neuron (control vs +ABO day 6)',
        png=f'{OUT_B1}/astro_neuron_4_3_26_figure.png',
        xlsx=f'{OUT_B1}/astro_neuron_AQuA2_summary.xlsx',
        exp_col='experiment', exp_id='4.3.26',
        treat='+ABO', ctrl='control',
        note='',
    ),
    dict(
        label='4/6/26 — Astro only (control vs +ABO)',
        png=f'{OUT_B2}/astro_neuron_4_6_26_figure.png',
        xlsx=f'{OUT_B2}/astro_neuron_batch2_AQuA2_summary.xlsx',
        exp_col='experiment', exp_id='4.6.26',
        treat='+ABO', ctrl='control',
        note='',
    ),
    dict(
        label='4/8/26 — Organoid (control vs +ABO)',
        png=f'{OUT_B2}/astro_neuron_4_8_26_figure.png',
        xlsx=f'{OUT_B2}/astro_neuron_batch2_AQuA2_summary.xlsx',
        exp_col='experiment', exp_id='4.8.26',
        treat='+ABO', ctrl='control',
        note='Very few events detected (0–3/sample)',
    ),
    dict(
        label='4/9/26 — GCaMP8s Asteroids d6 (control vs +amyloid)',
        png=f'{OUT_B2}/astro_neuron_4_9_26_figure.png',
        xlsx=f'{OUT_B2}/astro_neuron_batch2_AQuA2_summary.xlsx',
        exp_col='experiment', exp_id='4.9.26',
        treat='+amyloid', ctrl='control',
        note='High variance: rep3 amyloid=96 events',
    ),
    dict(
        label='4/10/26 — GCaMP8s Astroid d7 (control vs +amyloid)',
        png=f'{OUT_B2}/astro_neuron_4_10_26_figure.png',
        xlsx=f'{OUT_B2}/astro_neuron_batch2_AQuA2_summary.xlsx',
        exp_col='experiment', exp_id='4.10.26',
        treat='+amyloid', ctrl='control',
        note='',
    ),
]

# ── Colors ────────────────────────────────────────────────────────────────────
C_TEAL  = RGBColor(0x4D, 0xAD, 0xA0)
C_PINK  = RGBColor(0xE8, 0x74, 0x8A)
C_DARK  = RGBColor(0x2B, 0x2B, 0x2B)
C_GRAY  = RGBColor(0xF0, 0xF0, 0xF0)
C_WHITE = RGBColor(0xFF, 0xFF, 0xFF)
C_HEAD  = RGBColor(0x3A, 0x3A, 0x3A)


def set_cell(cell, text, bold=False, fontsize=8, align=PP_ALIGN.CENTER,
             bg=None, fg=C_DARK):
    cell.text = str(text)
    tf = cell.text_frame
    tf.word_wrap = False
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.runs[0] if p.runs else p.add_run()
    run.text = str(text)
    run.font.size = Pt(fontsize)
    run.font.bold = bold
    run.font.color.rgb = fg
    if bg:
        cell.fill.solid()
        cell.fill.fore_color.rgb = bg


def load_summary(exp):
    df = pd.read_excel(exp['xlsx'], sheet_name='Per-Sample Summary')
    if exp['exp_col'] and exp['exp_id']:
        df = df[df[exp['exp_col']] == exp['exp_id']].copy()
    df['rep'] = df['rep'].astype(str)
    return df


# ── Build presentation ────────────────────────────────────────────────────────
prs = Presentation()
prs.slide_width  = Inches(13.33)
prs.slide_height = Inches(7.5)

blank_layout = prs.slide_layouts[6]  # completely blank

# ── Title slide ───────────────────────────────────────────────────────────────
slide = prs.slides.add_slide(blank_layout)

# Background
bg = slide.shapes.add_shape(1, 0, 0, prs.slide_width, prs.slide_height)
bg.fill.solid(); bg.fill.fore_color.rgb = RGBColor(0x1A, 0x1A, 0x2E)
bg.line.fill.background()

tx = slide.shapes.add_textbox(Inches(1), Inches(2.2), Inches(11.3), Inches(1.5))
tf = tx.text_frame; tf.word_wrap = True
p = tf.paragraphs[0]
p.text = 'AQuA2 Ca²⁺ Event Analysis'
p.alignment = PP_ALIGN.CENTER
p.runs[0].font.size = Pt(36); p.runs[0].font.bold = True
p.runs[0].font.color.rgb = C_WHITE

p2 = tf.add_paragraph()
p2.text = 'Astrocyte + Neuron Experiments  ·  3/25/26 – 4/10/26'
p2.alignment = PP_ALIGN.CENTER
p2.runs[0].font.size = Pt(18)
p2.runs[0].font.color.rgb = RGBColor(0xAA, 0xCC, 0xCC)

tx2 = slide.shapes.add_textbox(Inches(1), Inches(4.2), Inches(11.3), Inches(0.6))
tf2 = tx2.text_frame
p3 = tf2.paragraphs[0]
p3.text = '8 experiments  ·  Panels: event timecourse (c), area vs duration scatter (d), event count (e), area (f), max ΔF/F (g)'
p3.alignment = PP_ALIGN.CENTER
p3.runs[0].font.size = Pt(12)
p3.runs[0].font.color.rgb = RGBColor(0x88, 0xAA, 0xAA)

# ── One slide per experiment ───────────────────────────────────────────────────
for exp in EXPERIMENTS:
    slide = prs.slides.add_slide(blank_layout)

    # Slide title bar
    title_bar = slide.shapes.add_shape(1, 0, 0, prs.slide_width, Inches(0.55))
    title_bar.fill.solid(); title_bar.fill.fore_color.rgb = C_HEAD
    title_bar.line.fill.background()

    tx = slide.shapes.add_textbox(Inches(0.15), Inches(0.05), Inches(10), Inches(0.45))
    tf = tx.text_frame
    p = tf.paragraphs[0]
    p.text = exp['label']
    p.runs[0].font.size = Pt(13); p.runs[0].font.bold = True
    p.runs[0].font.color.rgb = C_WHITE

    if exp['note']:
        nt = slide.shapes.add_textbox(Inches(10.0), Inches(0.08), Inches(3.2), Inches(0.4))
        tf_n = nt.text_frame
        p_n = tf_n.paragraphs[0]
        p_n.text = '⚠ ' + exp['note']
        p_n.alignment = PP_ALIGN.RIGHT
        p_n.runs[0].font.size = Pt(7.5)
        p_n.runs[0].font.color.rgb = RGBColor(0xFF, 0xCC, 0x66)

    # Figure image — full width, below title
    img_top  = Inches(0.58)
    img_h    = Inches(3.6)
    slide.shapes.add_picture(exp['png'], Inches(0), img_top,
                             width=prs.slide_width, height=img_h)

    # Data table
    df = load_summary(exp)
    cols = ['condition', 'rep', 'n_events', 'mean_area', 'mean_dff', 'mean_auc']
    cols = [c for c in cols if c in df.columns]
    headers = ['Condition', 'Rep', 'Events', 'Mean Area (px²)', 'Mean Max ΔF/F', 'Mean AUC']
    headers = headers[:len(cols)]

    n_rows = len(df) + 1  # +1 header
    tbl_top  = Inches(4.32)
    tbl_left = Inches(0.25)
    tbl_w    = Inches(12.83)
    tbl_h    = Inches(0.3 * n_rows + 0.05)

    table = slide.shapes.add_table(n_rows, len(cols), tbl_left, tbl_top,
                                   tbl_w, tbl_h).table

    # Column widths
    col_widths = [Inches(2.0), Inches(0.7), Inches(0.9), Inches(2.5), Inches(2.3), Inches(2.0)]
    for ci, w in enumerate(col_widths[:len(cols)]):
        table.columns[ci].width = w

    # Header row
    for ci, hdr in enumerate(headers):
        set_cell(table.cell(0, ci), hdr, bold=True, fontsize=8.5,
                 bg=C_HEAD, fg=C_WHITE)

    # Data rows
    treat = exp['treat']
    for ri, (_, row) in enumerate(df.iterrows(), start=1):
        is_treat = str(row.get('condition', '')) == treat
        row_bg = RGBColor(0xFF, 0xEE, 0xF0) if is_treat else RGBColor(0xEE, 0xF7, 0xF6)
        for ci, col in enumerate(cols):
            val = row[col]
            if col in ('mean_area', 'mean_dff', 'mean_auc') and pd.notna(val):
                val = f'{val:.1f}'
            elif pd.isna(val):
                val = '—'
            set_cell(table.cell(ri, ci), val, fontsize=8, bg=row_bg,
                     align=PP_ALIGN.CENTER if ci > 0 else PP_ALIGN.LEFT)

    # Legend
    leg = slide.shapes.add_textbox(Inches(0.25), Inches(7.1), Inches(12.83), Inches(0.3))
    tf_l = leg.text_frame
    p_l = tf_l.paragraphs[0]
    p_l.text = f'  ■ Control ({exp["ctrl"]})     ■ Treatment ({exp["treat"]})'
    run_l = p_l.runs[0]
    run_l.font.size = Pt(8)
    run_l.font.color.rgb = RGBColor(0x55, 0x55, 0x55)

# ── Combined summary slide ────────────────────────────────────────────────────
slide = prs.slides.add_slide(blank_layout)

title_bar = slide.shapes.add_shape(1, 0, 0, prs.slide_width, Inches(0.55))
title_bar.fill.solid(); title_bar.fill.fore_color.rgb = C_HEAD
title_bar.line.fill.background()
tx = slide.shapes.add_textbox(Inches(0.15), Inches(0.05), Inches(12), Inches(0.45))
tf = tx.text_frame
p = tf.paragraphs[0]
p.text = 'Summary — All Experiments'
p.runs[0].font.size = Pt(13); p.runs[0].font.bold = True
p.runs[0].font.color.rgb = C_WHITE

# Build combined table
rows_all = []
for exp in EXPERIMENTS:
    df = load_summary(exp)
    for cond in [exp['ctrl'], exp['treat']]:
        sub = df[df['condition'] == cond]
        if len(sub) == 0:
            continue
        n_ev   = sub['n_events'].sum()
        m_area = sub['mean_area'].mean()
        m_dff  = sub['mean_dff'].mean()
        rows_all.append({
            'Experiment': exp['label'].split(' — ')[0],
            'Condition': cond,
            'Type': 'Treatment' if cond == exp['treat'] else 'Control',
            'Total Events': int(n_ev),
            'Mean Area (px²)': f'{m_area:.0f}',
            'Mean Max ΔF/F': f'{m_dff:.2f}',
        })

df_sum = pd.DataFrame(rows_all)
n_rows_sum = len(df_sum) + 1
tbl = slide.shapes.add_table(n_rows_sum, 6,
                              Inches(0.25), Inches(0.65),
                              Inches(12.83), Inches(6.7)).table

sum_hdrs = ['Experiment', 'Condition', 'Type', 'Total Events', 'Mean Area (px²)', 'Mean Max ΔF/F']
sum_widths = [Inches(1.6), Inches(1.8), Inches(1.1), Inches(1.4), Inches(2.0), Inches(2.0)]
for ci, w in enumerate(sum_widths):
    tbl.columns[ci].width = w

for ci, hdr in enumerate(sum_hdrs):
    set_cell(tbl.cell(0, ci), hdr, bold=True, fontsize=8.5, bg=C_HEAD, fg=C_WHITE)

for ri, (_, row) in enumerate(df_sum.iterrows(), start=1):
    is_treat = row['Type'] == 'Treatment'
    row_bg = RGBColor(0xFF, 0xEE, 0xF0) if is_treat else RGBColor(0xEE, 0xF7, 0xF6)
    for ci, key in enumerate(sum_hdrs):
        align = PP_ALIGN.LEFT if ci < 2 else PP_ALIGN.CENTER
        set_cell(tbl.cell(ri, ci), row[key], fontsize=8, bg=row_bg, align=align)

# ── Save ──────────────────────────────────────────────────────────────────────
out_path = '/home/crnlqz/krenciklab/astrocyte+neuron_out_batch3/astro_neuron_AQuA2_all_experiments.pptx'
prs.save(out_path)
print(f'Saved: {out_path}')
