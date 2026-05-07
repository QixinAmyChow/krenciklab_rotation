"""
Histogram: Ca²⁺ event counts vs days post-ABO treatment.
Y = n_events per rep (mean ± SEM); X = day post-ABO; facet = cell culture type.
Appends a new slide to astro_neuron_AQuA2_v2.pptx.
"""

import numpy as np
import pandas as pd
import matplotlib
matplotlib.use('Agg')
import matplotlib.pyplot as plt
from scipy import stats
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN

OUT_B1   = '/home/crnlqz/krenciklab/astrocyte+neuron_out'
OUT_B2   = '/home/crnlqz/krenciklab/astrocyte+neuron_out_batch2'
OUT_B3   = '/home/crnlqz/krenciklab/astrocyte+neuron_out_batch3'
HIST_PNG = f'{OUT_B3}/abo_event_histogram.png'
PPTX_PATH = f'{OUT_B3}/astro_neuron_AQuA2_v2.pptx'

C_TREAT = '#E8748A'   # pink — +ABO / +amyloid
C_CTRL  = '#4DADA0'   # teal — control

# ── Facet / timepoint registry ────────────────────────────────────────────────
FACETS = [
    {
        'title': 'Astro only',
        'timepoints': [
            dict(label='day+3', exp_id='3.31.26',
                 xlsx=f'{OUT_B1}/astro_neuron_AQuA2_summary.xlsx',
                 treat='+amyloid', ctrl='control'),
            dict(label='day+3', exp_id='4.6.26',
                 xlsx=f'{OUT_B2}/astro_neuron_batch2_AQuA2_summary.xlsx',
                 treat='+ABO',    ctrl='control'),
        ],
    },
    {
        'title': 'Astro+Neuron',
        'timepoints': [
            dict(label='day+4', exp_id='4.1.26',
                 xlsx=f'{OUT_B1}/astro_neuron_AQuA2_summary.xlsx',
                 treat='+ABO', ctrl='control'),
            dict(label='day+6', exp_id='4.3.26',
                 xlsx=f'{OUT_B1}/astro_neuron_AQuA2_summary.xlsx',
                 treat='+ABO', ctrl='control'),
        ],
    },
    {
        'title': 'Organoid / Astroid',
        'timepoints': [
            dict(label='day+5', exp_id='4.8.26',
                 xlsx=f'{OUT_B2}/astro_neuron_batch2_AQuA2_summary.xlsx',
                 treat='+ABO',     ctrl='control'),
            dict(label='d6',    exp_id='4.9.26',
                 xlsx=f'{OUT_B2}/astro_neuron_batch2_AQuA2_summary.xlsx',
                 treat='+amyloid', ctrl='control'),
            dict(label='d7',    exp_id='4.10.26',
                 xlsx=f'{OUT_B2}/astro_neuron_batch2_AQuA2_summary.xlsx',
                 treat='+amyloid', ctrl='control'),
        ],
    },
]


def load_nevents(tp):
    df = pd.read_excel(tp['xlsx'], sheet_name='Per-Sample Summary')
    df = df[df['experiment'] == tp['exp_id']].copy()
    df['rep'] = df['rep'].astype(str)
    treat = df[df['condition'] == tp['treat']]['n_events'].values.astype(float)
    ctrl  = df[df['condition'] == tp['ctrl']]['n_events'].values.astype(float)
    return treat, ctrl


def pval_stars(a, b):
    if len(a) < 2 or len(b) < 2:
        return ''
    _, p = stats.ttest_ind(a, b)
    return '***' if p < 0.001 else '**' if p < 0.01 else '*' if p < 0.05 else 'n.s.'


# ── Figure ────────────────────────────────────────────────────────────────────
plt.rcParams.update({
    'font.size': 9, 'axes.spines.top': False, 'axes.spines.right': False,
    'axes.linewidth': 0.8, 'xtick.major.width': 0.8, 'ytick.major.width': 0.8,
})

fig, axes = plt.subplots(1, 3, figsize=(13, 4.2), sharey=True)
fig.suptitle('Ca²⁺ event counts vs days post-ABO/amyloid treatment\n(mean ± SEM; dots = individual reps)',
             fontsize=10, fontweight='bold', y=1.03)

BAR_W  = 0.32
JIT    = np.array([-0.07, 0.0, 0.07])

for ax, facet in zip(axes, FACETS):
    n_tp   = len(facet['timepoints'])
    x_pos  = np.arange(n_tp)          # one group per timepoint
    offset = {'treat': -BAR_W/2 - 0.02, 'ctrl': BAR_W/2 + 0.02}

    for xi, tp in enumerate(facet['timepoints']):
        treat_vals, ctrl_vals = load_nevents(tp)

        for vals, cond_key, color in [
            (treat_vals, 'treat', C_TREAT),
            (ctrl_vals,  'ctrl',  C_CTRL),
        ]:
            xc = xi + offset[cond_key]
            mean_v = vals.mean() if len(vals) else 0
            sem_v  = vals.std(ddof=1) / np.sqrt(len(vals)) if len(vals) > 1 else 0

            ax.bar(xc, mean_v, width=BAR_W, color=color, alpha=0.85,
                   yerr=sem_v, capsize=3,
                   error_kw=dict(lw=1.1, capthick=1.1), zorder=2)

            j = JIT[:len(vals)]
            ax.scatter(xc + j * 0.6, vals, color='white', edgecolors=color,
                       s=30, zorder=4, linewidths=1.1)

        # pair lines between treat and ctrl reps
        n_pairs = min(len(treat_vals), len(ctrl_vals))
        for k in range(n_pairs):
            ax.plot([xi + offset['treat'] + JIT[k] * 0.6,
                     xi + offset['ctrl']  + JIT[k] * 0.6],
                    [treat_vals[k], ctrl_vals[k]],
                    color='#888', lw=0.6, alpha=0.5, zorder=3)

        # significance annotation
        stars = pval_stars(treat_vals, ctrl_vals)
        if stars:
            ymax = max(treat_vals.max() if len(treat_vals) else 0,
                       ctrl_vals.max()  if len(ctrl_vals)  else 0)
            ytop = ymax * 1.12
            ax.plot([xi + offset['treat'], xi + offset['treat'],
                     xi + offset['ctrl'],  xi + offset['ctrl']],
                    [ytop * 0.94, ytop, ytop, ytop * 0.94], lw=0.8, color='k')
            ax.text(xi, ytop * 1.02, stars, ha='center', va='bottom', fontsize=8)

    ax.set_xticks(x_pos)
    ax.set_xticklabels([tp['label'] for tp in facet['timepoints']], fontsize=8.5)
    ax.set_title(facet['title'], fontsize=10, fontweight='bold')
    ax.set_xlabel('Days post treatment', fontsize=8)

axes[0].set_ylabel('Ca²⁺ events (n)', fontsize=9)

# shared legend
from matplotlib.patches import Patch
fig.legend(handles=[Patch(color=C_TREAT, label='+ABO / +amyloid'),
                    Patch(color=C_CTRL,  label='control')],
           loc='lower center', ncol=2, frameon=False, fontsize=9,
           bbox_to_anchor=(0.5, -0.07))


plt.tight_layout()
fig.savefig(HIST_PNG, dpi=180, bbox_inches='tight')
plt.close(fig)
print(f'Histogram saved: {HIST_PNG}')


# ── Append slide to PPTX ──────────────────────────────────────────────────────
prs = Presentation(PPTX_PATH)
blank = prs.slide_layouts[6]
slide = prs.slides.add_slide(blank)

# title bar
bar = slide.shapes.add_shape(1, 0, 0, prs.slide_width, Inches(0.55))
bar.fill.solid(); bar.fill.fore_color.rgb = RGBColor(0x1A, 0x23, 0x4A)
bar.line.fill.background()
tx = slide.shapes.add_textbox(Inches(0.15), Inches(0.06), Inches(13), Inches(0.43))
p = tx.text_frame.paragraphs[0]
p.text = 'Ca²⁺ Event Count vs Days Post-ABO/Amyloid Treatment  ·  Faceted by Cell Culture Type'
p.runs[0].font.size = Pt(12); p.runs[0].font.bold = True
p.runs[0].font.color.rgb = RGBColor(0xFF, 0xFF, 0xFF)

slide.shapes.add_picture(HIST_PNG, Inches(0.4), Inches(0.65),
                          width=Inches(12.53), height=Inches(6.7))

prs.save(PPTX_PATH)
print(f'Slide appended → {PPTX_PATH}  (now {len(prs.slides)} slides)')
